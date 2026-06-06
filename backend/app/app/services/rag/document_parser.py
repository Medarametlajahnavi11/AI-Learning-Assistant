from io import BytesIO

from docx import Document as DocxDocument
from pypdf import PdfReader
from pptx import Presentation


SUPPORTED_TYPES = {
    "application/pdf": "pdf",
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "docx",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation": "pptx",
    "text/plain": "txt",
}


def parse_document(content: bytes, content_type: str) -> str:
    if content_type not in SUPPORTED_TYPES:
        raise ValueError("Unsupported file type")

    file_kind = SUPPORTED_TYPES[content_type]

    if file_kind == "pdf":
        reader = PdfReader(BytesIO(content))
        return "\n".join(page.extract_text() or "" for page in reader.pages)

    if file_kind == "docx":
        doc = DocxDocument(BytesIO(content))
        return "\n".join(p.text for p in doc.paragraphs)

    if file_kind == "pptx":
        prs = Presentation(BytesIO(content))
        lines: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                text = getattr(shape, "text", None)
                if text:
                    lines.append(text)
        return "\n".join(lines)

    return content.decode("utf-8", errors="ignore")
